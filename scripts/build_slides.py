"""Generate the slide deck: livrables/presentation_sujet9_bases_vectorielles.pptx.

Self-supporting slides — every talking point appears on the slide as a concise
bullet (keywords, not sentences) so the speaker can present from the slides alone.
Theme: navy + emerald with kicker chips, 14-slide pedagogical flow.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "livrables" / "presentation_sujet9_bases_vectorielles.pptx"

FONT = "Calibri"
EMU_PER_PX = 9525
TOTAL = 14

C = {
    "bg": "ffffff", "ink": "0b2545", "text": "243b53", "muted": "6b7a90",
    "line": "e2e8f0", "card": "f4f7fb",
    "teal": "0e7c66", "coral": "d9480f", "violet": "5b3fa8", "amber": "b7791f", "navy": "11366b",
    "tealtint": "e3f2ee", "coraltint": "fbe9e0", "violettint": "ece7f7",
    "ambertint": "f7efdd", "navytint": "e6ecf6",
}
TINT = {"teal": "tealtint", "coral": "coraltint", "violet": "violettint",
        "amber": "ambertint", "navy": "navytint", "ink": "navytint"}


def px(v: float) -> Emu:
    return Emu(int(round(v * EMU_PER_PX)))


def rgb(name: str) -> RGBColor:
    return RGBColor.from_string(C.get(name, name).upper())


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb("bg")
    return s


def shape(slide, left, top, w, h, fill="none", line="none", radius=0.0):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, px(left), px(top), px(w), px(h))
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill == "none":
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    if line == "none":
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line); shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(slide, txt, left, top, w, h, size=17, bold=False, color="text",
         align="left", anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(px(left), px(top), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    amap = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    for i, line in enumerate(str(txt).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = amap[align]
        r = p.add_run(); r.text = line
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = rgb(color)
    return box


def chip(slide, label, left, top, accent):
    w = len(label) * 9.2 + 38
    shape(slide, left, top, w, 30, TINT[accent], "none", radius=0.5)
    text(slide, label.upper(), left + 14, top, w - 22, 30, 11, True, accent,
         anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    return w


def header(slide, kicker, index, accent="teal"):
    chip(slide, kicker, 64, 40, accent)
    text(slide, f"{index}/{TOTAL}", 1140, 42, 76, 22, 12, True, "muted", "right")


def title(slide, ttl, subtitle, index, kicker="Sujet 9", accent="teal"):
    header(slide, kicker, index, accent)
    text(slide, ttl, 64, 80, 1150, 96, 30, True, "ink")
    if subtitle:
        text(slide, subtitle, 64, 168, 1150, 40, 17, False, "muted")


def bullets(slide, items, x, y, w, gap=54, size=17, dot="teal"):
    for i, it in enumerate(items):
        top = y + i * gap
        bold = isinstance(it, tuple)
        label = it[0] if bold else it
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(top + 8), px(8), px(8))
        d.fill.solid(); d.fill.fore_color.rgb = rgb(dot); d.line.fill.background()
        d.shadow.inherit = False
        # bold lead segment before " : " rendered by writing two runs
        box = slide.shapes.add_textbox(px(x + 20), px(top), px(w - 20), px(gap - 2))
        tf = box.text_frame; tf.word_wrap = True
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]
        if " — " in label:
            lead, rest = label.split(" — ", 1)
            r1 = p.add_run(); r1.text = lead + " — "
            r1.font.name = FONT; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = rgb("ink")
            r2 = p.add_run(); r2.text = rest
            r2.font.name = FONT; r2.font.size = Pt(size); r2.font.color.rgb = rgb("text")
        else:
            r = p.add_run(); r.text = label
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = rgb("text")


def callout(slide, left, top, w, h, body, accent="teal"):
    shape(slide, left, top, w, h, TINT[accent], accent, radius=0.08)
    text(slide, body, left + 18, top + 14, w - 36, h - 24, 16, True, "ink", anchor=MSO_ANCHOR.MIDDLE)


def metric(slide, left, top, w, h, value, label, accent="teal"):
    shape(slide, left, top, w, h, "card", "line", 0.08)
    text(slide, value, left + 14, top + 16, w - 28, 46, 30, True, accent, "center")
    text(slide, label, left + 14, top + 68, w - 28, h - 76, 13.5, False, "text", "center")


def image_contain(slide, rel, left, top, bw, bh):
    p = ROOT / rel
    with PILImage.open(p) as im:
        iw, ih = im.size
    ratio = min(bw / iw, bh / ih)
    w, h = iw * ratio, ih * ratio
    slide.shapes.add_picture(str(p), px(left + (bw - w) / 2), px(top + (bh - h) / 2), px(w), px(h))


def notes(slide, lines):
    slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def build():
    prs = Presentation()
    prs.slide_width = px(1280); prs.slide_height = px(720)

    # 1 — Titre
    s = new_slide(prs)
    shape(s, 64, 92, 7, 150, "teal", "none", 0.4)
    text(s, "Bases vectorielles et recherche sémantique", 92, 84, 1010, 100, 40, True, "ink")
    text(s, "Retrouver les bons documents quand la requête n'a pas les mêmes mots — et à grande échelle",
         92, 196, 1040, 36, 18, False, "muted")
    chip(s, "Sujet 9", 92, 256, "teal")
    bullets(s, [
        "Q1 — l'hybride BM25 + re-ranking bat-il BM25 seul et le dense seul ?",
        "Q2 — à partir de quelle taille HNSW devient-il indispensable, et à quel coût ?",
    ], 92, 308, 1040, gap=46, size=18, dot="violet")
    metric(s, 92, 410, 230, 120, "2 000", "documents", "violet")
    metric(s, 346, 410, 230, 120, "20", "requêtes", "amber")
    metric(s, 600, 410, 230, 120, "4", "régimes comparés", "teal")
    text(s, "Nemo MULLER et Yoan VIOLLET", 92, 566, 600, 30, 18, True, "ink")
    text(s, "M1 Algorithmique et ingénierie de données", 92, 598, 600, 26, 14, False, "muted")
    notes(s, ["Présenter le sujet et les deux questions qui guident tout l'exposé."])

    # 2 — Requête, document, résultat
    s = new_slide(prs)
    title(s, "Requête, document, résultat", "Trois notions à distinguer avant les algorithmes.", 2, "Les bases", "violet")
    bullets(s, [
        "Requête — ce que l'utilisateur tape au moment de chercher",
        "Document — un texte déjà stocké dans la base",
        "Résultat — un document affiché après comparaison requête ↔ base",
        "Le moteur compare la requête à toute la base, puis classe",
    ], 64, 248, 700, gap=66, size=18, dot="violet")
    callout(s, 64, 542, 700, 80,
            "À retenir : tous les résultats sont des documents, mais tous les documents ne deviennent pas des résultats.", "violet")
    # mini schéma à droite
    shape(s, 820, 250, 360, 300, "card", "line", 0.05)
    text(s, "Requête", 850, 280, 300, 26, 16, True, "ink")
    shape(s, 850, 312, 300, 40, "violettint", "violet", 0.2)
    text(s, "moteur : compare + classe", 862, 322, 276, 22, 13, False, "violet")
    text(s, "↓", 985, 356, 30, 26, 18, True, "muted")
    text(s, "Documents (toute la base)", 850, 386, 300, 24, 14, True, "muted")
    shape(s, 850, 414, 300, 110, "navytint", "line", 0.08)
    shape(s, 880, 446, 240, 48, "tealtint", "teal", 0.15)
    text(s, "Résultats (top-K)", 905, 460, 200, 22, 13, True, "teal")
    notes(s, ["Définir les 3 mots. Phrase-ancre : résultats ⊂ documents."])

    # 3 — Le problème
    s = new_slide(prs)
    title(s, "Le problème : les bons mots ne coïncident pas", "Une requête et un document peuvent parler de la même chose sans les mêmes mots.", 3, "Le problème", "coral")
    callout(s, 64, 240, 1150, 70,
            "Exemple : « repérer un paiement anormal »  ↔  « détection de transactions frauduleuses »", "coral")
    bullets(s, [
        "Recherche lexicale — cherche les mots communs → rate si le vocabulaire diffère",
        "Recherche sémantique — rapproche les textes par le sens",
        "Enjeu RAG — si la recherche rate le bon passage, la réponse générée est fausse",
    ], 64, 350, 1150, gap=66, size=18, dot="coral")
    notes(s, ["Exemple concret de paraphrase ; introduire les deux familles de méthodes."])

    # 4 — BM25
    s = new_slide(prs)
    title(s, "BM25 — la recherche lexicale", "Score selon les mots partagés entre requête et document.", 4, "Méthode 1", "teal")
    bullets(s, [
        "Index inversé — chaque mot pointe vers les documents qui le contiennent",
        "IDF — un mot rare (« HNSW ») pèse plus qu'un mot fréquent (« le »)",
        "k1 — sature la répétition d'un mot ; b — normalise par la longueur",
        "Implémenté from scratch — formule entièrement explicable",
    ], 64, 246, 720, gap=62, size=17, dot="teal")
    callout(s, 820, 250, 360, 130, "Force : excellent quand les mots exacts sont présents et discriminants.", "teal")
    callout(s, 820, 398, 360, 130, "Limite : aveugle au sens — synonymes et paraphrases sans mot commun.", "coral")
    notes(s, ["Index inversé → IDF → k1/b. Insister sur le from scratch. Force/limite."])

    # 5 — Embeddings
    s = new_slide(prs)
    title(s, "Embeddings — la recherche sémantique", "Le texte devient un vecteur ; le sens proche = des vecteurs proches.", 5, "Méthode 2", "violet")
    image_contain(s, "outputs/intro_embeddings/cosine_similarity_heatmap.png", 64, 244, 500, 380)
    bullets(s, [
        "Texte → vecteur dense — modèle all-MiniLM-L6-v2, 384 dimensions",
        "Requête et documents encodés dans le même espace",
        "Comparaison — similarité cosinus (angle) ; normalisé → produit scalaire",
        "Sens proche → direction proche dans l'espace",
    ], 600, 256, 600, gap=64, size=17, dot="violet")
    text(s, "Heatmap (20 phrases) : les phrases d'un même thème ont une similarité plus élevée.",
         600, 540, 600, 50, 14, False, "muted")
    notes(s, ["384 dimensions, cosinus = angle. La heatmap rend l'idée visible."])

    # 6 — Recherche exacte et échelle
    s = new_slide(prs)
    title(s, "Recherche exacte et passage à l'échelle", "Comparer la requête à tous les vecteurs devient lourd quand le corpus grandit.", 6, "Le défi", "amber")
    image_contain(s, "outputs/figures/distance_concentration.png", 64, 250, 560, 370)
    bullets(s, [
        "Recherche exacte — compare la requête à TOUS les vecteurs",
        "Coût O(n·d) — linéaire, lourd à grande échelle",
        "Haute dimension — les distances se concentrent (CV 0,505 → 0,023)",
        "KD-tree / Ball Tree — éliminent moins de régions → inefficaces",
        "→ on passe aux méthodes approchées (ANN)",
    ], 660, 250, 540, gap=58, size=16.5, dot="amber")
    notes(s, ["Exact = O(n·d). La dispersion des distances s'effondre avec la dimension."])

    # 7 — HNSW
    s = new_slide(prs)
    title(s, "HNSW — un index approché par graphe", "On accepte une approximation contrôlée pour aller bien plus vite.", 7, "La solution ANN", "teal")
    image_contain(s, "outputs/figures/hnsw_layers.png", 64, 248, 580, 372)
    bullets(s, [
        "Graphe multi-couches — couches hautes = raccourcis (peu de points)",
        "Couche basse — contient tous les vecteurs",
        "Recherche — BFS guidé par la distance (voisins prometteurs)",
        "efSearch — nombre de candidats explorés (vitesse ↔ qualité)",
        "Qualité mesurée — recall@K vs recherche exacte",
    ], 680, 250, 520, gap=58, size=16.5, dot="teal")
    notes(s, ["Couches = raccourcis ; descente vers la zone pertinente ; efSearch règle l'effort."])

    # 8 — Notre projet
    s = new_slide(prs)
    title(s, "Notre projet — un banc d'essai reproductible", "On contrôle tout pour pouvoir mesurer proprement.", 8, "Le projet", "navy")
    metric(s, 64, 246, 250, 116, "2 000", "documents synthétiques", "violet")
    metric(s, 340, 246, 250, 116, "20", "requêtes annotées", "amber")
    metric(s, 616, 246, 250, 116, "4", "régimes comparés", "teal")
    metric(s, 892, 246, 250, 116, "42", "graine (reproductible)", "navy")
    bullets(s, [
        "Corpus synthétique — on connaît la vérité terrain (docs pertinents par requête)",
        "Régimes comparés — BM25, dense, hybride, et HNSW pour l'indexation",
        "Métriques calculables et résultats régénérables à l'identique",
    ], 64, 408, 1150, gap=58, size=17, dot="navy")
    notes(s, ["Le synthétique donne une vérité terrain exacte → métriques fiables."])

    # 9 — Pipeline hybride
    s = new_slide(prs)
    title(s, "Le pipeline hybride — filtrer puis réordonner", "Deux étages : rappel (BM25) puis précision (re-ranking).", 9, "Méthode 3", "violet")
    # flow
    steps = [("1. BM25", "top-50 candidats\n(rappel large, rapide)", "teal"),
             ("2. Re-ranking", "réordonne les 50\n(bi- ou cross-encodeur)", "violet"),
             ("3. Top-10", "classement final", "amber")]
    for i, (t, b, a) in enumerate(steps):
        left = 64 + i * 396
        shape(s, left, 252, 348, 150, "card", "line", 0.06)
        shape(s, left, 252, 6, 150, a, "none", 0.4)
        text(s, t, left + 20, 270, 312, 30, 19, True, "ink")
        text(s, b, left + 20, 306, 312, 80, 16, False, "text")
        if i < 2:
            text(s, "→", left + 356, 300, 36, 44, 26, True, "violet", "center")
    callout(s, 64, 432, 1150, 110,
            "Point clé : le re-ranker ne voit QUE les 50 candidats de BM25. Un bon document absent du top-50 est perdu — ce plafond explique nos résultats.", "coral")
    text(s, "Re-ranker de référence : cross-encoder ms-marco-MiniLM-L-6-v2.", 64, 566, 1150, 26, 14, False, "muted")
    notes(s, ["Deux étages ; le plafond de rappel du 1er étage prépare le résultat hybride."])

    # 10 — Évaluation
    s = new_slide(prs)
    title(s, "Comment on mesure la qualité", "Une mesure qui discrimine, sinon les comparaisons ne veulent rien dire.", 10, "Évaluation", "amber")
    bullets(s, [
        "Problème — une pertinence binaire trop large saturait les scores",
        "Solution — pertinence graduée 0-3",
        "Grade 3 — bon thème + bon concept principal + bon domaine",
        "Métriques — precision@10, recall@10, nDCG@10, MAP@50",
        "nDCG — tient compte de l'ordre (bon doc en 1er > en 10e)",
    ], 64, 246, 720, gap=58, size=17, dot="amber")
    metric(s, 832, 250, 168, 120, "3", "min / requête", "coral")
    metric(s, 1016, 250, 168, 120, "11,5", "médiane", "violet")
    metric(s, 924, 388, 168, 120, "20", "max / requête", "teal")
    notes(s, ["Binaire saturait → graduée 0-3, gold ~5-30 docs/requête. nDCG = ordre."])

    # 11 — Résultat hybride
    s = new_slide(prs)
    title(s, "Résultat 1 — l'hybride : nuancé", "Bat BM25, fait jeu égal avec le dense, gagne sur le sémantique.", 11, "Résultats", "teal")
    image_contain(s, "outputs/figures/regime_comparison.png", 64, 248, 640, 372)
    bullets(s, [
        "nDCG@10 — BM25 0,540 · dense 0,617 · hybride 0,614",
        "Hybride > BM25, mais jeu égal avec le dense (écart 0,003 = bruit)",
        "Pourquoi — plafonné par le 1er étage BM25 (rappel borné au top-50)",
        "Sémantique — hybride 0,508 > dense 0,478 > BM25 0,377",
    ], 730, 250, 470, gap=58, size=16, dot="teal")
    callout(s, 730, 542, 470, 82, "Valeur de l'hybride : pas un gain partout, une robustesse ciblée quand lexical et sémantique se complètent.", "teal")
    notes(s, ["Honnête : jeu égal avec le dense ; gain réel sur les paraphrases ; plafond BM25."])

    # 12 — Résultat HNSW
    s = new_slide(prs)
    title(s, "Résultat 2 — HNSW passe à l'échelle", "Inutile petit, indispensable vers 50 000 vecteurs.", 12, "Résultats", "coral")
    image_contain(s, "outputs/figures/scale_crossover.png", 64, 248, 660, 372)
    bullets(s, [
        "Petite taille — HNSW peu utile (exact déjà instantané)",
        "Bascule pratique — ~50 000 vecteurs (gain ≥ 5×)",
        "1 million — ~48× plus rapide que la force brute",
        "Coût — recall 0,995 → 0,75 à efSearch fixe ; +18 % mémoire",
        "Récupérable — augmenter efSearch (au prix de la latence)",
    ], 750, 250, 460, gap=56, size=16, dot="coral")
    notes(s, ["Bascule pratique ~50k ; 48× à 1M ; coût en recall réglable par efSearch."])

    # 13 — Démo
    s = new_slide(prs)
    title(s, "Démonstration — Streamlit", "Rendre les résultats visibles, en direct.", 13, "Démo", "violet")
    image_contain(s, "outputs/figures/scale_recall_latency.png", 64, 250, 600, 370)
    bullets(s, [
        "Comparaison BM25 / dense / hybride sur une requête sémantique",
        "Re-ranking — quels documents montent / descendent dans le top-10",
        "Slider efSearch — recall ↑ et latence ↑ (le compromis en direct)",
        "Plan B — courbes pré-calculées si la démo ne se lance pas",
    ], 700, 256, 500, gap=64, size=17, dot="violet")
    notes(s, ["Une requête sémantique, le re-ranking, puis le slider efSearch. Plan B = figures."])

    # 14 — Conclusion
    s = new_slide(prs)
    title(s, "Conclusion — pas de méthode magique", "Chaque approche a son domaine de pertinence.", 14, "Conclusion", "teal")
    cards = [("BM25", "mots exacts, termes rares", "teal"),
             ("Embeddings", "le sens, les paraphrases", "violet"),
             ("Hybride", "surtout les requêtes sémantiques", "amber"),
             ("HNSW", "le passage à l'échelle", "coral")]
    for i, (t, b, a) in enumerate(cards):
        left = 64 + i * 290
        shape(s, left, 246, 268, 150, "card", "line", 0.06)
        shape(s, left, 246, 6, 150, a, "none", 0.4)
        text(s, t, left + 20, 266, 230, 30, 19, True, "ink")
        text(s, b, left + 20, 304, 230, 80, 16, False, "text")
    callout(s, 64, 424, 558, 130,
            "Limite : corpus synthétique → on défend des tendances, pas des magnitudes absolues.", "coral")
    callout(s, 656, 424, 558, 130,
            "Ouverture : fusion BM25 + dense avant re-ranking ; pgvector + filtres SQL ; vrais embeddings à 1M.", "teal")
    text(s, "Notre apport : montrer dans quels cas chaque méthode devient pertinente — pas qu'une seule gagne.",
         64, 576, 1150, 36, 16, True, "ink", "center")
    notes(s, ["Répondre aux 2 questions sans surpromettre ; assumer les limites ; ouvrir."])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
